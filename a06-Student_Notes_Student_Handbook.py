#!/usr/bin/env python3
"""
Speaker Notes Generation Module

This module handles the generation of speaker notes for PowerPoint slides
using batch processing with an LLM (similar to how image prompts are generated).
"""

import json
import os
import sys
import time
from pptx import Presentation

# Add the directory to the path so we can import modules from the project
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import LLM utilities
from a06_Image_Generation import call_llm

# Constants
OUTPUT_PPTX = "course_presentation.pptx"
OUTLINE_FILE = "course_outline.txt"

# Get directory from current_directory.txt if available
def get_current_directory():
    try:
        if os.path.exists("current_directory.txt"):
            with open("current_directory.txt", "r") as dir_file:
                directory = dir_file.read().strip()
                if directory:
                    print(f"Using directory from current_directory.txt: {directory}")
                    # Create directory if it doesn't exist
                    if not os.path.exists(directory):
                        os.makedirs(directory)
                        print(f"Created directory: {directory}")
                    return directory
    except Exception as e:
        print(f"Error reading current_directory.txt: {str(e)}")
    return None

# Current project directory
CURRENT_DIR = get_current_directory()

# Global variable to store generated notes
slide_notes = {}

def generate_slides_info_from_outline(outline_data):
    """
    Extract slide information from the outline data structure
    
    Args:
        outline_data: The parsed outline data from the course_outline.txt
        
    Returns:
        List of dictionaries containing slide title and content
    """
    slides_info = []
    
    # First add the title slide
    if CURRENT_DIR:
        course_name = os.path.basename(CURRENT_DIR).replace('_', ' ')
        slides_info.append({
            "title": f"Course: {course_name}",
            "content": "Welcome to this comprehensive course on " + course_name,
            "type": "title"
        })
        print(f"Added title slide for course: {course_name}")
    
    # Process modules, topics, and subtopics to collect slide content
    for module, module_data in outline_data.items():
        # Extract module number without the word "Module"
        module_number = module.split()[1]  # Gets the number after "Module"
        module_title = f"{module_number}: {module_data['title']}"
        
        # Module slide content
        topic_bullets = [f"{topic_data['title']}" for topic, topic_data in module_data["topics"].items()]
        module_content = "\n".join(topic_bullets)
        slides_info.append({"title": module_title, "content": module_content, "type": "module"})
        print(f"Added module slide: {module_title}")
        
        # Process topics
        for topic, topic_data in module_data["topics"].items():
            # Extract topic number 
            topic_parts = topic.split()
            topic_number = topic_parts[1]  # Gets the number like "1.1"
            topic_title = f"{topic_number}: {topic_data['title']}"
            
            # Topic slide content
            subtopic_bullets = [f"{subtopic_data['title']}" for subtopic, subtopic_data in topic_data["subtopics"].items()]
            topic_content = "\n".join(subtopic_bullets)
            slides_info.append({"title": topic_title, "content": topic_content, "type": "topic"})
            
            # Process subtopics
            for subtopic, subtopic_data in topic_data["subtopics"].items():
                subtopic_title = subtopic_data['title']
                subtopic_content = "\n".join(subtopic_data["points"])
                slides_info.append({"title": subtopic_title, "content": subtopic_content, "type": "subtopic"})
    
    return slides_info

def generate_all_speaker_notes(slides_info, max_slides=0):
    """
    Generate speaker notes for all slides in the presentation in a single batch
    
    Args:
        slides_info: List of dictionaries containing slide information (title, content, type)
        max_slides: Maximum number of slides to generate notes for (0 = no limit)
        
    Returns:
        Dictionary of slide titles to speaker notes
    """
    global slide_notes
    
    print("\nPreparing slide content for batch speaker notes generation...")
    
    # Verify we have slide info to process
    if not slides_info:
        print("No slide information provided for notes generation.")
        # Check if we should use a fallback
        if CURRENT_DIR:
            # Extract course name from directory path
            dir_parts = CURRENT_DIR.split('/')
            if len(dir_parts) >= 2:
                course_name = dir_parts[-1]
                print(f"Using fallback with course name: {course_name}")
                
                # Add a fallback slide for notes generation
                slides_info.append({
                    "title": f"Course: {course_name}",
                    "content": "Introduction to the course overview and fundamentals.",
                    "type": "title"
                })
                print("Added fallback slide for notes generation")
        
        # If still empty, return empty dict
        if not slides_info:
            return {}
        
    print(f"Processing {len(slides_info)} slides for notes generation")
    
    # For debugging
    for slide in slides_info:
        print(f"Slide: {slide['title']}")
        content_preview = slide['content'][:50] + '...' if len(slide['content']) > 50 else slide['content']
        print(f"Content: {content_preview}")
        print("---")
    
    print(f"Found {len(slides_info)} slides for speaker notes generation")
    
    # Limit slides if needed
    if max_slides > 0 and max_slides < len(slides_info):
        slides_info = slides_info[:max_slides]
        print(f"Limited to {len(slides_info)} slides for testing")
    
    # Create a single batch prompt for all slides
    system_prompt = """
    You are an expert educator and speaker notes writer for technical presentations.
    
    I will provide you with a list of slide titles and content for a technical course.
    For EACH slide, create comprehensive speaker notes that would help an instructor deliver the content effectively.
    
    PAY SPECIAL ATTENTION to these specific slide types:
    1. For slides with "type": "title", create a welcoming introduction for the entire course.
    2. For slides with "type": "module", create an engaging module introduction.
    3. For all other slides, create standard detailed notes.
    
    IMPORTANT: EVERY slide must have notes, including title slides and module slides.
    
    The student notes should:
    - Be a couple of paragraphs long and can be read by a student to understand the content
    - For title slides: Include a warm welcome and brief course overview
    - For module slides: Include an introduction to the module's importance and key learning objectives
    
    Provide your response as a JSON object with the following structure:
    {"slide_title": "student notes"}
    
    For example:
    {
      "Course: AI in Cybersecurity": "Welcome to our comprehensive course on AI in Cybersecurity. Throughout this program, we'll explore how artificial intelligence is revolutionizing both offensive and defensive cybersecurity operations. This course will equip you with practical knowledge of AI-powered security tools and strategies for implementing them in your organization.",
      "1: Introduction to AI in Cybersecurity": "Welcome to our module on AI in Cybersecurity. This module provides an overview of how artificial intelligence technologies are transforming the cybersecurity landscape. Emphasize to participants that AI is both creating new security challenges and offering powerful new defensive capabilities. Engage the audience by asking how many of them currently use AI tools in their security operations. Transition: Let's begin by examining how cyber threats have evolved alongside AI technologies."
    }
    
    Only include the JSON response with no additional explanations or text.
    """
    
    # Build the user prompt with all slides
    slides_json = json.dumps(slides_info, indent=2)
    user_prompt = f"Generate detailed speaker notes for each of these technical presentation slides:\n{slides_json}"
    
    # Call the LLM
    try:
        print(f"Sending batch request to LLM for {len(slides_info)} slide notes...")
        speaker_notes_json = call_llm(user_prompt, system_prompt)
        
        # Save the raw notes response to a file
        output_filename = "06_Enhanced_Notes.txt"
        
        # Check if we should save to a subdirectory
        if CURRENT_DIR:
            # Ensure directory exists
            if not os.path.exists(CURRENT_DIR):
                os.makedirs(CURRENT_DIR)
            output_path = os.path.join(CURRENT_DIR, output_filename)
        else:
            output_path = output_filename
        
        # Save the raw response
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(speaker_notes_json)
        print(f"Enhanced notes saved to: {output_path}")
        
        # Parse the JSON response
        try:
            # Extract JSON from response, handling markdown code blocks if present
            if "```json" in speaker_notes_json:
                # Extract content between ```json and ```
                json_start = speaker_notes_json.find("```json") + 7
                json_end = speaker_notes_json.find("```", json_start)
                if json_end > json_start:
                    speaker_notes_json = speaker_notes_json[json_start:json_end].strip()
            
            # Try to clean up any extra whitespace or non-JSON content
            speaker_notes_json = speaker_notes_json.strip()
            if not speaker_notes_json.startswith("{"):
                # Find first occurrence of '{'
                start_idx = speaker_notes_json.find('{')
                if start_idx >= 0:
                    speaker_notes_json = speaker_notes_json[start_idx:]
            
            # Load the JSON
            slide_notes = json.loads(speaker_notes_json)
            print(f"Successfully received {len(slide_notes)} speaker notes from LLM")
            
        except json.JSONDecodeError as e:
            print(f"Failed to parse JSON response: {str(e)}")
            print("Raw response:", speaker_notes_json)
            # Use fallback notes if JSON parsing fails
            slide_notes = {}
            for slide_info in slides_info:
                slide_title = slide_info["title"]
                slide_type = slide_info["type"]
                # Create simple fallback notes
                if slide_type == "module":
                    slide_notes[slide_title] = f"Welcome to {slide_title}. This module covers key concepts and principles related to this topic. Review the bullet points on the slide and elaborate on each topic."
                elif slide_type == "topic":
                    slide_notes[slide_title] = f"This slide covers {slide_title}. Explain each subtopic listed and how they relate to the overall topic."
                else:  # subtopic
                    slide_notes[slide_title] = f"In this slide about {slide_title}, discuss each bullet point in detail, providing examples where appropriate."
    
    except Exception as e:
        print(f"Error generating speaker notes: {str(e)}")
        # Create simple fallback notes
        slide_notes = {}
        for slide_info in slides_info:
            slide_title = slide_info["title"]
            slide_notes[slide_title] = f"Discuss the key points on this slide about {slide_title}."
    
    return slide_notes


def get_speaker_notes(slide_title):
    """
    Retrieves speaker notes for a given slide title
    
    Args:
        slide_title: The title of the slide
        
    Returns:
        Speaker notes text or a fallback message
    """
    global slide_notes
    
    if slide_title in slide_notes:
        return slide_notes[slide_title]
    else:
        return f"Speaker notes for '{slide_title}' not available. Please explain the slide content based on the bullet points shown."


def add_speaker_notes_to_presentation(pptx_file, slide_notes_dict):
    """
    Add generated speaker notes to each slide in a PowerPoint presentation
    
    Args:
        pptx_file: Path to the PowerPoint file
        slide_notes_dict: Dictionary of slide titles to speaker notes
        
    Returns:
        True if successful, False otherwise
    """
    try:
        # Load the presentation
        presentation = Presentation(pptx_file)
        
        # Track how many notes we added
        notes_added = 0
        
        # Special handling for presentations with just a single slide (title slide)
        if len(presentation.slides) == 1:
            print("Adding notes to title slide in single-slide presentation")
            slide = presentation.slides[0]
            
            # Extract course name from current directory if available
            fallback_title = None
            if CURRENT_DIR:
                dir_parts = CURRENT_DIR.split('/')
                if len(dir_parts) >= 2:
                    course_name = dir_parts[-1]
                    fallback_title = f"Course: {course_name}"
            
            # Try to find notes for this title slide using various approaches
            notes_text = None
            
            # First try: Look for exact title match from any text in the slide
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    try:
                        slide_title = shape.text_frame.text
                        if slide_title and slide_title in slide_notes_dict:
                            notes_text = slide_notes_dict[slide_title]
                            print(f"Found notes using exact title match: {slide_title}")
                            break
                    except:
                        pass
            
            # Second try: Use fallback title from directory name
            if not notes_text and fallback_title and fallback_title in slide_notes_dict:
                notes_text = slide_notes_dict[fallback_title]
                print(f"Found notes using fallback title: {fallback_title}")
            
            # Third try: Just use the first note in the dictionary if we only have one
            if not notes_text and len(slide_notes_dict) == 1:
                first_key = next(iter(slide_notes_dict))
                notes_text = slide_notes_dict[first_key]
                print(f"Using the only available note for title slide")
            
            # If we found notes, add them to the slide
            if notes_text:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                
                # Clear existing notes if any
                if notes_text_frame.text:
                    for paragraph in notes_text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                
                # Add the new notes
                p = notes_text_frame.paragraphs[0] if notes_text_frame.paragraphs else notes_text_frame.add_paragraph()
                p.text = notes_text
                notes_added += 1
        
        # Regular handling for multi-slide presentations
        else:
            # Process each slide
            for slide in presentation.slides:
                # Find the title shape in the slide
                title_shape = None
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        # Check if this shape is the title (usually the first text shape)
                        if hasattr(shape, "name") and ("Title" in shape.name or notes_added == 0):
                            title_shape = shape
                            break
                
                # If we found a title, try to add notes
                if title_shape and title_shape.has_text_frame:
                    slide_title = title_shape.text_frame.text
                    
                    # Look for notes for this title
                    if slide_title in slide_notes_dict:
                        # Get the notes slide
                        notes_slide = slide.notes_slide
                        
                        # Clear any existing notes text (if needed)
                        notes_text_frame = notes_slide.notes_text_frame
                        if notes_text_frame.text:
                            for paragraph in notes_text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = ""
                        
                        # Add the new notes
                        notes_text = slide_notes_dict[slide_title]
                        p = notes_text_frame.paragraphs[0] if notes_text_frame.paragraphs else notes_text_frame.add_paragraph()
                        p.text = notes_text
                        notes_added += 1
        
        # Save the presentation
        presentation.save(pptx_file)
        
        print(f"Successfully added speaker notes to {notes_added} slides in '{pptx_file}'")
        return True
        
    except Exception as e:
        print(f"Error adding speaker notes to presentation: {str(e)}")
        return False


def process_presentation_with_notes(outline_data=None, pptx_file=None, max_slides=0):
    """
    Main function to generate speaker notes and add them to a PowerPoint presentation
    
    Args:
        outline_data: The parsed outline data structure (can be None if extracting directly from PPT)
        pptx_file: Path to the PowerPoint presentation file
        max_slides: Maximum number of slides to process (0 = no limit)
        
    Returns:
        True if successful, False otherwise
    """
    # Use the specified file or determine from current directory
    if pptx_file is None:
        if CURRENT_DIR:
            pptx_file = os.path.join(CURRENT_DIR, OUTPUT_PPTX)
        else:
            pptx_file = OUTPUT_PPTX
    
    # Check if the presentation exists
    if not os.path.exists(pptx_file):
        print(f"Error: PowerPoint file '{pptx_file}' not found.")
        print("Please run a05_CREATE_POWERPOINT.py first to generate the presentation.")
        return False
        
    # If outline_data is None, extract slides directly from the PowerPoint
    if outline_data is None or not outline_data:
        print("No outline data provided. Extracting slide content directly from PowerPoint...")
        try:
            # Open the presentation
            presentation = Presentation(pptx_file)
            print(f"PowerPoint file opened: {pptx_file}")
            print(f"Total slides in presentation: {len(presentation.slides)}")
            
            # Process slides to extract content
            slides_info = []
            
            # Get course title from directory name if available
            course_title = ""
            if CURRENT_DIR:
                base_dir = os.path.basename(CURRENT_DIR)
                course_title = base_dir.replace('_', ' ')
            
            # Special handling for the first slide - always treat as a title slide
            print("Processing title slide")
            if len(presentation.slides) > 0:
                slide = presentation.slides[0]
                
                # Try to extract presentation title from the properties or first slide
                try:
                    # First try from core properties
                    ppt_title = presentation.core_properties.title
                    
                    # If not found, try to extract from slide content
                    if not ppt_title and hasattr(slide, "shapes"):
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape.text_frame and hasattr(shape.text_frame, "text"):
                                if shape.text_frame.text:
                                    ppt_title = shape.text_frame.text
                                    break
                    
                    # If still not found, use course title from directory
                    if not ppt_title and course_title:
                        ppt_title = f"Course: {course_title}"
                    
                    # If we have a title, add the slide
                    if ppt_title:
                        print(f"Using title slide with title: {ppt_title}")
                        slides_info.append({
                            "title": ppt_title, 
                            "content": "Welcome to this comprehensive course. This title slide introduces the main topic and sets expectations for the learning journey ahead.",
                            "type": "title"
                        })
                        print("Added title slide for notes generation")
                    else:
                        # No title found, create a generic one
                        if course_title:
                            slides_info.append({
                                "title": f"Course: {course_title}",
                                "content": "Welcome slide for the course",
                                "type": "title"
                            })
                            print(f"Added generic title slide for {course_title}")
                except Exception as e:
                    print(f"Error processing title slide: {str(e)}")
                    # Still add a generic title slide
                    if course_title:
                        slides_info.append({
                            "title": f"Course: {course_title}",
                            "content": "Welcome slide for the course",
                            "type": "title"
                        })
                        print(f"Added generic title slide for {course_title}")
            
            # Regular slide processing
            for i, slide in enumerate(presentation.slides):
                print(f"Processing slide {i+1}/{len(presentation.slides)}")
                
                slide_title = ""
                slide_content = []
                
                # List all shapes in the slide
                print(f"  Shapes in slide {i+1}: {len(slide.shapes)}")
                for j, shape in enumerate(slide.shapes):
                    shape_name = getattr(shape, 'name', 'Unknown')
                    print(f"  Shape {j+1}: Type={type(shape).__name__}, Name={shape_name}")
                    
                    # Try to get text content if available
                    if hasattr(shape, "text_frame"):
                        try:
                            # Access text frames and paragraphs directly
                            if shape.text_frame and hasattr(shape.text_frame, "paragraphs"):
                                frame_text = ""
                                for para in shape.text_frame.paragraphs:
                                    if para.text:
                                        frame_text += para.text + "\n"
                                
                                # Remove trailing newline
                                frame_text = frame_text.strip()
                                
                                if frame_text:
                                    print(f"    Text: {frame_text[:50]}{'...' if len(frame_text) > 50 else ''}")
                                    # If this is the title shape or first slide and has "Title" in name
                                    if hasattr(shape, "name") and ("Title" in shape.name or i == 0):
                                        slide_title = frame_text
                                        print(f"    Identified as slide title")
                                    # Otherwise treat as content
                                    else:
                                        slide_content.append(frame_text)
                                        print(f"    Added as slide content")
                            # Fallback to direct text property
                            elif hasattr(shape.text_frame, "text") and shape.text_frame.text:
                                text = shape.text_frame.text
                                print(f"    Text: {text[:50]}{'...' if len(text) > 50 else ''}")
                                if hasattr(shape, "name") and ("Title" in shape.name or i == 0):
                                    slide_title = text
                                    print(f"    Identified as slide title")
                                else:
                                    slide_content.append(text)
                                    print(f"    Added as slide content")
                        except Exception as e:
                            print(f"    Error extracting text: {str(e)}")
                            import traceback
                            traceback.print_exc()
                
                # Try to find a title if we don't have one yet - use first text shape if needed
                if not slide_title and slide_content:
                    slide_title = slide_content[0]
                    slide_content = slide_content[1:] if len(slide_content) > 1 else []
                    print(f"  Using first text element as title: {slide_title[:50]}{'...' if len(slide_title) > 50 else ''}")
                
                # Determine slide type
                slide_type = "slide"  # Default type
                
                # Title slide detection
                if i == 0:
                    slide_type = "title"
                    print(f"  Detected as title slide: {slide_title}")
                
                # Module slide detection
                elif "Module" in slide_title or ("Module" in slide_title.lower() and ":" in slide_title):
                    slide_type = "module"
                    print(f"  Detected as module slide: {slide_title}")
                
                # Only add slides that have a title
                if slide_title:
                    print(f"  Adding slide with title: {slide_title} (Type: {slide_type})")
                    slides_info.append({
                        "title": slide_title,
                        "content": "\n".join(slide_content),
                        "type": slide_type
                    })
                else:
                    # Even for slides without title, create a generic one based on slide number
                    generic_title = f"Slide {i+1}"
                    print(f"  Creating generic title for slide {i+1}: {generic_title}")
                    slides_info.append({
                        "title": generic_title,
                        "content": "\n".join(slide_content) if slide_content else f"Content for slide {i+1}",
                        "type": slide_type
                    })
            
            print(f"Extracted content from {len(slides_info)} slides")
            
        except Exception as e:
            print(f"Error extracting slide content: {str(e)}")
            return False
    else:
        # Use the provided outline data to generate slides info
        slides_info = generate_slides_info_from_outline(outline_data)
    
    # Generate speaker notes
    all_notes = generate_all_speaker_notes(slides_info, max_slides)
    
    if not all_notes:
        print("Failed to generate speaker notes.")
        return False
    
    # Add notes to the presentation
    return add_speaker_notes_to_presentation(pptx_file, all_notes)


if __name__ == "__main__":
    from a05_CREATE_POWERPOINT import parse_outline
    
    # Determine paths based on current directory
    outline_path = OUTLINE_FILE
    output_path = OUTPUT_PPTX
    
    if CURRENT_DIR:
        # Use directory-specific paths
        outline_path = os.path.join(CURRENT_DIR, OUTLINE_FILE)
        output_path = os.path.join(CURRENT_DIR, OUTPUT_PPTX)
    
    # Check if the outline file exists
    if not os.path.exists(outline_path):
        print(f"\nError: Course outline file '{outline_path}' not found.")
        print("Please run a04_CREATE_OUTLINE.py first to generate the course outline.")
        sys.exit(1)
    
    # Check if the presentation file exists
    if not os.path.exists(output_path):
        print(f"\nError: PowerPoint file '{output_path}' not found.")
        print("Please run a05_CREATE_POWERPOINT.py first to generate the presentation.")
        sys.exit(1)
    
    # Parse the outline file
    try:
        outline_data, course_title = parse_outline(outline_path)
        
        print(f"\nProcessing presentation '{output_path}' for student notes...")
        
        # Get the course title
        course_title = ""
        if CURRENT_DIR:
            course_title = os.path.basename(CURRENT_DIR).replace('_', ' ')
        elif course_title:
            # Already have course title from parse_outline
            pass
        else:
            course_title = "This Course"
        
        print(f"Course title: {course_title}")
        
        # Process the presentation with speaker notes (no slide limit)
        success = process_presentation_with_notes(outline_data, output_path)
        
        if success:
            print(f"\nStudent notes added to presentation '{output_path}'.")
            print("="*70)
        else:
            print(f"\nFailed to add student notes to presentation '{output_path}'.")
        
    except Exception as e:
        print(f"\nError occurred: {str(e)}")
        sys.exit(1)
