#!/usr/bin/env python3
"""
PowerPoint Generator from Course Outline

This script reads a course outline text file and generates a PowerPoint presentation
with slides organized by modules, topics, and subtopics. Each slide includes an AI-generated
image based on the slide content.
"""

import time
import os
import sys
import re
import json
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import concurrent.futures

# Import our custom slide snapshot generator
from a07_Slide_Snapshot_Generator import generate_snapshots_for_presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import random
import math
from PIL import Image, ImageDraw, ImageFont

# Import image generation and speaker notes modules
from a06_Image_Generation import generate_image_for_slide, get_enhanced_prompt, generate_all_image_prompts
from arunware_image_generator import generate_images_parallel

# Import speaker notes generation module
try:
    from a06_Student_Notes_Student_Handbook import process_presentation_with_notes
    SPEAKER_NOTES_AVAILABLE = True
except ImportError:
    print("Speaker notes module not available.")
    SPEAKER_NOTES_AVAILABLE = False

# Import the LLM access module
import importlib.util
import sys

# Import the LLM module dynamically (renamed file with 'a' prefix)
llm_spec = importlib.util.spec_from_file_location("llm_module", "a02_LLM_Access.py")
llm_module = importlib.util.module_from_spec(llm_spec)
sys.modules["llm_module"] = llm_module
llm_spec.loader.exec_module(llm_module)
call_llm = llm_module.call_llm

# Constants for enhanced notes
ENHANCED_NOTES_FILE = "06_Enhanced_Notes.txt"

# Function to load enhanced notes
def load_enhanced_notes(notes_file):
    """
    Load enhanced speaker notes from the notes file.
    
    Args:
        notes_file: Path to the enhanced notes file
        
    Returns:
        Dictionary mapping slide titles/identifiers to speaker notes
    """
    print(f"Loading enhanced speaker notes from {notes_file}...")
    
    try:
        with open(notes_file, 'r') as f:
            content = f.read()
            # Check if the file is in JSON format
            try:
                notes_data = json.loads(content)
                print(f"Loaded {len(notes_data)} enhanced speaker notes entries.")
                return notes_data
            except json.JSONDecodeError:
                # If it's not valid JSON, try to extract JSON from the content
                json_match = re.search(r'```json\n(.+?)\n```', content, re.DOTALL)
                if json_match:
                    notes_data = json.loads(json_match.group(1))
                    print(f"Extracted and loaded {len(notes_data)} enhanced speaker notes entries.")
                    return notes_data
                else:
                    print("Error: Could not parse enhanced speaker notes as JSON.")
                    return {}
    except FileNotFoundError:
        print(f"Enhanced speaker notes file not found at {notes_file}")
        return {}
    except Exception as e:
        print(f"Error loading enhanced speaker notes: {str(e)}")
        return {}

# Define constants
OUTLINE_FILE = "course_outline.txt"
OUTPUT_PPTX = "course_presentation.pptx"
OUTPUT_MARKDOWN = "course_presentation.md"
# Try both template formats
TEMPLATE_OPTIONS = ["template.pptx", "template.potx"]

# Get directory from current_directory.txt if available
def get_current_directory():
    try:
        if os.path.exists("current_directory.txt"):
            with open("current_directory.txt", "r") as dir_file:
                directory = dir_file.read().strip()
                if directory:
                    print(f"Using directory from current_directory.txt: {directory}")
                    # Create directory if it doesn't exist
                    if not os.path.exists(directory):
                        os.makedirs(directory)
                        print(f"Created directory: {directory}")
                    return directory
    except Exception as e:
        print(f"Error reading current_directory.txt: {str(e)}")
    return None

# Current project directory
CURRENT_DIR = get_current_directory()

# Constants for slide generation

# Testing limit
MAX_SLIDES_TO_PROCESS = 50

# Function to reset and manage the slide counter
def reset_slide_counter():
    global SLIDES_PROCESSED
    SLIDES_PROCESSED = 0
    
# Initialize counter
SLIDES_PROCESSED = 0  # Counter to track total slides generated

# Define PowerPoint generation constants

# Define colors
COLORS = {
    "module_title": RGBColor(0, 32, 96),     # Dark blue
    "topic_title": RGBColor(0, 112, 192),    # Medium blue
    "subtopic_title": RGBColor(0, 176, 240), # Light blue
    "bullet_text": RGBColor(0, 0, 0),         # Black
    "background": RGBColor(255, 255, 255),  # White
    "title": RGBColor(0, 68, 129),          # Cisco Blue
    "accent1": RGBColor(0, 155, 229),       # Light blue
    "accent2": RGBColor(100, 195, 84),      # Green
    "accent3": RGBColor(206, 59, 50)        # Red/orange
}

# Constants for title slide image generation
TITLE_IMAGE_WIDTH = 1600
TITLE_IMAGE_HEIGHT = 900  # 16:9 aspect ratio

def generate_title_slide_image(title_text, output_dir, subtitle_text="Generated Course Presentation"):
    """Generate a professional title slide image based on the course title.
    
    Args:
        title_text: The main course title text
        output_dir: Directory to save the image
        subtitle_text: Optional subtitle text
        
    Returns:
        Path to the generated title slide image
    """
    # Create the output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Path for the generated title image
    title_image_path = os.path.join(output_dir, "title_slide_image.png")
    
    # Create a blank image with a gradient background
    img = Image.new('RGB', (TITLE_IMAGE_WIDTH, TITLE_IMAGE_HEIGHT), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # Create a nice gradient background (from dark blue to lighter blue)
    for y in range(TITLE_IMAGE_HEIGHT):
        # Calculate gradient color (from dark blue to light blue)
        progress = y / TITLE_IMAGE_HEIGHT
        r = int(0 + (0 * progress))
        g = int(32 + (123 * progress))
        b = int(96 + (133 * progress))
        
        # Draw a horizontal line with this color
        draw.line([(0, y), (TITLE_IMAGE_WIDTH, y)], fill=(r, g, b))
    
    # Try to load a nice font, fallback to default if not available
    title_font_size = 80
    subtitle_font_size = 40
    
    try:
        # Try to find a nice font on the system
        system_fonts = [
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Supplemental/Georgia.ttf"
        ]
        
        font_path = None
        for font in system_fonts:
            if os.path.exists(font):
                font_path = font
                break
        
        if font_path:
            title_font = ImageFont.truetype(font_path, title_font_size)
            subtitle_font = ImageFont.truetype(font_path, subtitle_font_size)
        else:
            # Fallback to default
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
            
    except Exception as e:
        print(f"Error loading font: {e}")
        # Fallback to default
        title_font = ImageFont.load_default()
        subtitle_font = ImageFont.load_default()
    
    # Add a decorative element (abstract shapes)
    # Draw some circles in the background
    for _ in range(10):
        # Random position and size
        x = random.randint(0, TITLE_IMAGE_WIDTH)
        y = random.randint(0, TITLE_IMAGE_HEIGHT)
        size = random.randint(100, 400)
        
        # Semi-transparent white
        opacity = random.randint(30, 80)
        circle_color = (255, 255, 255, opacity)
        
        # Draw circle
        draw.ellipse((x-size//2, y-size//2, x+size//2, y+size//2), 
                     fill=None, outline=circle_color, width=3)
    
    # Add title text (centered horizontally, positioned in upper half)
    title_y = TITLE_IMAGE_HEIGHT // 3
    
    # Wrap title text if it's too long
    max_width = TITLE_IMAGE_WIDTH - 200  # Leave some margin
    words = title_text.split()
    lines = []
    current_line = []
    
    # Simple text wrapping logic
    for word in words:
        test_line = ' '.join(current_line + [word])
        # Use getbbox to get text size
        if hasattr(title_font, 'getbbox'):
            text_width = title_font.getbbox(test_line)[2]
        else:
            # Fallback method for older PIL versions
            text_width = title_font.getmask(test_line).getbbox()[2]
            
        if text_width <= max_width:
            current_line.append(word)
        else:
            lines.append(' '.join(current_line))
            current_line = [word]
    
    if current_line:
        lines.append(' '.join(current_line))
    
    # Draw each line of the title
    for i, line in enumerate(lines):
        # Use getbbox to get text size
        if hasattr(title_font, 'getbbox'):
            text_width = title_font.getbbox(line)[2]
        else:
            # Fallback method for older PIL versions
            text_width = title_font.getmask(line).getbbox()[2]
            
        text_x = (TITLE_IMAGE_WIDTH - text_width) // 2
        # Add white text with a slight offset for shadow effect
        draw.text((text_x+2, title_y+2+i*title_font_size), line, font=title_font, fill=(0, 0, 0, 100))
        draw.text((text_x, title_y+i*title_font_size), line, font=title_font, fill=(255, 255, 255))
    
    # Add subtitle
    subtitle_y = title_y + (len(lines) * title_font_size) + 40
    
    # Use getbbox to get text size
    if hasattr(subtitle_font, 'getbbox'):
        subtitle_width = subtitle_font.getbbox(subtitle_text)[2]
    else:
        # Fallback method for older PIL versions
        subtitle_width = subtitle_font.getmask(subtitle_text).getbbox()[2]
        
    subtitle_x = (TITLE_IMAGE_WIDTH - subtitle_width) // 2
    
    # Add subtitle text with shadow
    draw.text((subtitle_x+1, subtitle_y+1), subtitle_text, font=subtitle_font, fill=(0, 0, 0, 100))
    draw.text((subtitle_x, subtitle_y), subtitle_text, font=subtitle_font, fill=(255, 255, 255))
    
    # Add a decorative line under the subtitle
    line_y = subtitle_y + subtitle_font_size + 10
    line_width = TITLE_IMAGE_WIDTH // 3
    line_x_start = (TITLE_IMAGE_WIDTH - line_width) // 2
    line_x_end = line_x_start + line_width
    draw.line([(line_x_start, line_y), (line_x_end, line_y)], fill=(255, 255, 255), width=3)
    
    # Save the image
    try:
        # Convert to RGB mode if it has alpha channel
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        img.save(title_image_path)
        print(f"Generated title slide image: {title_image_path}")
        return title_image_path
    except Exception as e:
        print(f"Error saving title slide image: {e}")
        return None

def create_title_slide(prs, title_text, image_path=None):
    """Create the title slide for the presentation using Title Slide Layout.
    
    Args:
        prs: PowerPoint presentation object
        title_text: Title text for the slide
        image_path: Optional path to an image to include on the title slide
    """
    # Extract course name from title_text
    course_name = title_text
    # Find the Title Slide Layout (typically index 0)
    title_slide_layout = None
    
    # Try to find the layout by name first
    for layout in prs.slide_layouts:
        if layout.name == 'Title Slide' or 'title slide' in layout.name.lower():
            title_slide_layout = layout
            break
    
    # If not found by name, use the first layout (index 0)
    if not title_slide_layout:
        title_slide_layout = prs.slide_layouts[0]
    
    print(f"Using Layout '{title_slide_layout.name}' for title slide")
    
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Find the title and subtitle placeholders
    title = None
    subtitle = None
    
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # Title placeholder
            title = shape
        elif shape.placeholder_format.type == 2:  # Subtitle placeholder
            subtitle = shape
    
    # Set title and subtitle text
    if title:
        print(f"Found title placeholder, setting text to: {title_text}")
        title.text = title_text
        # Format title
        title_format = title.text_frame.paragraphs[0].font
        title_format.size = Pt(44)
        title_format.bold = True
        title_format.color.rgb = COLORS["module_title"]
    
    if subtitle:
        subtitle.text = "Generated Course Presentation"
    
    # Add image if provided
    if image_path and os.path.exists(image_path):
        # Calculate image position and size
        # For a title slide, we'll make the image fill the entire slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Image will take up the entire slide
        left = 0  # Start at left edge
        top = 0   # Start at top edge
        width = slide_width  # Full slide width
        height = slide_height  # Full slide height
        
        # Add the image first (before title text) so it becomes the background
        print(f"Adding full-slide image to title slide: {image_path}")
        try:
            # Add image to cover the entire slide
            picture = slide.shapes.add_picture(image_path, left, top, width, height)
            
            # Move the picture to the back so text appears on top
            picture.z_order = 0  # Put at the back of the z-order
            
            # If the image is now the background, make sure title text is visible
            if title:
                # Update title text formatting for better visibility on image background
                title_frame = title.text_frame
                for paragraph in title_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    paragraph.font.bold = True
                    
                    # Add shadow to make text more readable over image
                    paragraph.font.shadow.inherit = False
                    paragraph.font.shadow.visible = True
            
            # Same for subtitle if it exists
            if subtitle:
                subtitle_frame = subtitle.text_frame
                for paragraph in subtitle_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(240, 240, 240)  # Light gray
                    paragraph.font.bold = True
                    paragraph.font.shadow.inherit = False
                    paragraph.font.shadow.visible = True
        except Exception as e:
            print(f"Error adding image to title slide: {e}")
    
    # Add speaker notes to the title slide from enhanced notes if available
    try:
        # Try to get notes from global enhanced_notes dictionary
        found_notes = False
        if 'enhanced_notes' in globals() and enhanced_notes:
            # Try variations of the title to find a match in the enhanced notes
            course_variants = [
                course_name,
                f"Course: {course_name}", 
                "Course: Cisco AI Defense",  # Special case from what we observed
                f"Title: {course_name}",
                "Title Slide"
            ]
            
            # Also try variations without "Fundamentals" if present
            base_name = course_name.replace(" Fundamentals", "")
            if base_name != course_name:
                course_variants.append(base_name)
                course_variants.append(f"Course: {base_name}")
            
            # Try all variant keys
            for key in course_variants:
                if key in enhanced_notes:
                    welcome_notes = enhanced_notes[key]
                    slide.notes_slide.notes_text_frame.text = welcome_notes
                    print(f"Added enhanced speaker notes to title slide using key: '{key}'")
                    found_notes = True
                    break
            
            # Special handling: try getting the first key in enhanced_notes if it looks like a title
            if not found_notes and enhanced_notes:
                for key in enhanced_notes.keys():
                    if key.lower().startswith("course:") or "title" in key.lower():
                        welcome_notes = enhanced_notes[key]
                        slide.notes_slide.notes_text_frame.text = welcome_notes
                        print(f"Added enhanced speaker notes to title slide using key: '{key}'")
                        found_notes = True
                        break
            
            # Still not found, use fallback
            if not found_notes:
                # Generate fallback notes
                welcome_notes = f"Welcome to the course: {course_name}. This comprehensive course will provide you with in-depth knowledge and practical skills about {course_name}. Throughout this course, you will learn key concepts, best practices, and hands-on techniques that you can apply in real-world scenarios. Let's begin our journey into {course_name}."
                slide.notes_slide.notes_text_frame.text = welcome_notes
                print(f"Added fallback speaker notes to title slide (enhanced notes not found)")
                print(f"Available keys were: {list(enhanced_notes.keys())}")
        else:  # No enhanced notes available
            # Generate fallback notes
            welcome_notes = f"Welcome to the course: {course_name}. This comprehensive course will provide you with in-depth knowledge and practical skills about {course_name}. Throughout this course, you will learn key concepts, best practices, and hands-on techniques that you can apply in real-world scenarios. Let's begin our journey into {course_name}."
            slide.notes_slide.notes_text_frame.text = welcome_notes
            print(f"Added fallback speaker notes to title slide (enhanced notes not available)")
    except Exception as e:
        print(f"Error adding notes to title slide: {e}")
    
    return slide

def create_content_slide(prs, title_text, bullet_points, title_color=None, generate_images=True, image_path=None):
    """Create a content slide with title and bullet points."""
    # Use the content slide layout
    slide_layout = prs.slide_layouts[1]  # Layout index 1 is for Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # If no title color specified, use default color
    if title_color is None:
        title_color = COLORS.get("title", RGBColor(0, 68, 129))  # Default to Cisco blue
    
    # Set the title
    title = slide.shapes.title
    title.text = title_text
    all_slide_text = ""
    all_slide_text += title_text
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(40)
    title_format.bold = True
    title_format.color.rgb = title_color
    
    # Add bullet points if there are any
    if bullet_points:
        content = slide.placeholders[1]  # Index 1 is the content placeholder
        text_frame = content.text_frame
        
        for i, point in enumerate(bullet_points):
            # First paragraph already exists in empty text frame
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
                
            paragraph.text = point
            all_slide_text += "\n" + point
            paragraph.level = 0  # Top level bullet
            
            # Format bullet text
            font = paragraph.font
            font.size = Pt(28)
            font.color.rgb = COLORS["bullet_text"]
    
    # Use provided image_path if available, otherwise generate if enabled
    selected_image = None
    
    if image_path and os.path.exists(image_path):
        # Use the pre-generated image
        selected_image = image_path
        print(f"Using pre-generated image for slide: {title_text}")
    elif generate_images:
        print(f"Generating new image for slide: {title_text}")
        try:
            # Generate images for this slide
            image_paths = generate_image_for_slide(title_text, all_slide_text)
            
            # If images were generated successfully
            if image_paths and isinstance(image_paths, list) and len(image_paths) > 0:
                # For now, use the first image for the slide
                selected_image = image_paths[0]
                print(f"Using image 1/{len(image_paths)} for slide: {selected_image}")
                print(f"Other images available in slide_images directory")
        except Exception as e:
            print(f"Error generating image: {str(e)}")
            print("Continuing without image for this slide...")
            selected_image = None
    
    # Add the image to the slide if we have one
    if selected_image:
        try:
            # Get slide dimensions (widescreen slide is 13.33 x 7.5 inches)
            slide_width = Inches(13.33)
            slide_height = Inches(7.5)
            
            # Calculate image width in inches while preserving aspect ratio
            # Original was 2.5 inches, making it 25% larger as requested
            image_width_inches = Inches(3.125)  # 2.5 * 1.25 = 3.125 (25% larger)
            
            # Calculate height based on aspect ratio (512:1024 = 1:2)
            image_height_inches = image_width_inches * 2
            
            # Position at bottom right with some margin
            left = slide_width - image_width_inches - Inches(0.1)  # Small margin from right
            top = slide_height - image_height_inches - Inches(0.1)  # Small margin from bottom
            
            # Add the selected image to the slide
            slide.shapes.add_picture(selected_image, left, top, width=image_width_inches)
            print(f"Added image to slide: {title_text}")
        except Exception as e:
            print(f"Error adding image to slide: {str(e)}")
    else:
        print(f"No image available for slide: {title_text}")
        
    # Add speaker notes for all content slides, with special handling for module slides
    try:
        slide_key = title_text  # Use the slide title as the key for the notes
        
        # Clean up module name for better matching if this is a module slide
        if "module" in title_text.lower():
            module_name = title_text.replace("Module", "").strip()
            if ":" in module_name:
                module_name = module_name.split(":")[1].strip()
        
        # Try to get notes from global enhanced_notes dictionary
        if 'enhanced_notes' in globals() and enhanced_notes:
            # Try variations of the title for better matching
            possible_keys = [title_text, f"{title_text}", title_text.strip()]
            
            for key in possible_keys:
                if key in enhanced_notes:
                    slide_notes = enhanced_notes[key]
                    slide.notes_slide.notes_text_frame.text = slide_notes
                    print(f"Added enhanced speaker notes to slide: {title_text}")
                    break
            else:  # No match found in the loop
                if "module" in title_text.lower():
                    # Generate fallback module notes
                    module_notes = f"Module: {module_name}. This module will introduce you to the key concepts and principles of {module_name}. We will explore the fundamental aspects, practical applications, and best practices related to this module topic. By the end of this module, you will have a solid understanding of {module_name} and be able to apply these concepts in various scenarios."
                    slide.notes_slide.notes_text_frame.text = module_notes
                    print(f"Added fallback speaker notes to module slide: {title_text} (enhanced notes not found)")
        elif "module" in title_text.lower():  # No enhanced notes available, but it's a module slide
            # Generate fallback module notes
            module_notes = f"Module: {module_name}. This module will introduce you to the key concepts and principles of {module_name}. We will explore the fundamental aspects, practical applications, and best practices related to this module topic. By the end of this module, you will have a solid understanding of {module_name} and be able to apply these concepts in various scenarios."
            slide.notes_slide.notes_text_frame.text = module_notes
            print(f"Added fallback speaker notes to module slide: {title_text} (enhanced notes not available)")
    except Exception as e:
        print(f"Error adding notes to slide: {e}")

    return slide

def parse_outline(file_path):
    """
    Parse the course outline text file and extract modules, topics, subtopics, and points.
    
    Returns a nested dictionary structure and the course title.
    """
    course_title = "Course Presentation"  # Default title
    
    with open(file_path, 'r') as f:
        lines = f.readlines()
    
    # First line is the course title
    if lines and not lines[0].strip().startswith("1."):
        course_title = lines[0].strip()
        lines = lines[1:]  # Remove the title line
    
    # Initialize the structure
    outline = {}
    current_module = None
    current_topic = None
    current_subtopic = None
    
    # Process line by line
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check for different levels based on the number of parts in the numbering
        parts = line.split(' ', 1)
        if len(parts) < 2:
            continue  # Skip lines without proper formatting
            
        number_part, title_part = parts
        
        # Check the format of the number to determine the level
        num_dots = number_part.count('.')
        
        # Remove trailing dot if present
        number_part = number_part.rstrip('.')
        
        if num_dots == 0 and number_part.isdigit():
            # It's a module (e.g., "1 Introduction to AI Security")
            module_number = number_part
            current_module = f"Module {module_number}"
            outline[current_module] = {
                "title": title_part.strip(),
                "topics": {}
            }
            current_topic = None
            current_subtopic = None
            print(f"Found module: {current_module} - {title_part.strip()}")
            
        elif num_dots == 1:
            # It's a topic (e.g., "1.1 AI Security Overview")
            if not current_module:
                # If no module has been defined yet, create a default one
                current_module = "Module 1"
                outline[current_module] = {
                    "title": "Main Module",
                    "topics": {}
                }
                print(f"Created default module: {current_module}")
                
            topic_number = number_part
            current_topic = f"Topic {topic_number}"
            outline[current_module]["topics"][current_topic] = {
                "title": title_part.strip(),
                "subtopics": {}
            }
            current_subtopic = None
            print(f"Found topic: {current_topic} - {title_part.strip()}")
            
        elif num_dots == 2:
            # It's a subtopic (e.g., "1.1.1 Defining AI Security")
            if not current_module:
                # If no module has been defined yet, create a default one
                current_module = "Module 1"
                outline[current_module] = {
                    "title": "Main Module",
                    "topics": {}
                }
                
            if not current_topic:
                # If no topic has been defined yet, create a default one
                topic_prefix = number_part.split('.')[0]
                current_topic = f"Topic {topic_prefix}.1"
                outline[current_module]["topics"][current_topic] = {
                    "title": "Main Topic",
                    "subtopics": {}
                }
                
            subtopic_number = number_part
            current_subtopic = f"Subtopic {subtopic_number}"
            outline[current_module]["topics"][current_topic]["subtopics"][current_subtopic] = {
                "title": title_part.strip(),
                "points": []
            }
            print(f"Found subtopic: {current_subtopic} - {title_part.strip()}")
            
        elif num_dots == 3:
            # It's a point (e.g., "1.1.1.1 Key AI Security Principle")
            if not current_module or not current_topic or not current_subtopic:
                # Skip if we don't have a proper hierarchy
                continue
                
            point = title_part.strip()
            outline[current_module]["topics"][current_topic]["subtopics"][current_subtopic]["points"].append(point)
            
        # Handle the case where we have modules without the "Module" prefix
        elif num_dots == 0 and re.match(r'^\d+$', number_part):
            # This looks like a top-level item (module)
            module_number = number_part
            current_module = f"Module {module_number}"
            outline[current_module] = {
                "title": title_part.strip(),
                "topics": {}
            }
            current_topic = None
            current_subtopic = None
            print(f"Found simple module: {current_module} - {title_part.strip()}")
    
    # Print summary of what was found
    module_count = len(outline)
    topic_count = sum(len(module_data["topics"]) for module_data in outline.values())
    subtopic_count = sum(
        sum(len(topic_data["subtopics"]) for topic_data in module_data["topics"].values())
        for module_data in outline.values()
    )
    
    print(f"Found {module_count} modules, {topic_count} topics, and {subtopic_count} subtopics in outline.")
    
    return outline, course_title

def prepare_slides_data(outline_data, max_slides=10):
    """
    Prepare data for slide generation without creating slides yet.
    This allows us to generate images in parallel before creating the slides.
    
    Args:
        outline_data: A nested dictionary containing the outline structure
        max_slides: Maximum number of slides to process
        
    Returns:
        List of slide data dictionaries containing title, content, and color information
    """
    slides_data = []
    
    # First slide will be title slide, so we don't include it here
    
    # Process modules
    for module, module_data in outline_data.items():
        if not module_data.get("topics"):
            continue
            
        # Module slide
        topic_bullets = [f"{topic_data['title']}" for topic, topic_data in module_data["topics"].items()]
        slides_data.append({
            "title": f"{module}: {module_data['title']}",
            "content": topic_bullets,
            "color": COLORS["module_title"],
            "type": "module"
        })
        
        if len(slides_data) >= max_slides:
            break
            
        # Process topics
        for topic, topic_data in module_data["topics"].items():
            if len(slides_data) >= max_slides:
                break
                
            # Extract topic number 
            topic_parts = topic.split()
            topic_number = topic_parts[1]  # Gets the number like "1.1"
            
            # Topic slide
            subtopic_bullets = [f"{subtopic_data['title']}" for subtopic, subtopic_data in topic_data["subtopics"].items()]
            slides_data.append({
                "title": f"{topic_number}: {topic_data['title']}",
                "content": subtopic_bullets,
                "color": COLORS["topic_title"],
                "type": "topic"
            })
            
            if len(slides_data) >= max_slides:
                break
                
            # Process subtopics
            for subtopic, subtopic_data in topic_data["subtopics"].items():
                if len(slides_data) >= max_slides:
                    break
                    
                # Subtopic slide
                slides_data.append({
                    "title": f"{subtopic_data['title']}",
                    "content": subtopic_data["points"],
                    "color": COLORS["subtopic_title"],
                    "type": "subtopic"
                })
                
                if len(slides_data) >= max_slides:
                    break
        
        # If we've reached the limit, break out of modules loop
        if len(slides_data) >= max_slides:
            break
    
    return slides_data


def generate_slide_images_parallel(slides_data, batch_size=25):
    """
    Generate images for multiple slides in parallel using the enhanced batch prompt generation.
    
    Args:
        slides_data: List of dictionaries with slide title and content
        batch_size: Number of images to generate in parallel
        
    Returns:
        Dictionary mapping slide titles to image file paths
    """
    from a06_Image_Generation import get_enhanced_prompt, IMAGE_DIR, IMAGE_WIDTH, IMAGE_HEIGHT
    
    print(f"\nGenerating images for {len(slides_data)} slides...")
    
    # First, generate all enhanced prompts in a single batch using LLM
    print("Step 1: Generating enhanced image prompts with LLM...")
    start_time = time.time()
    
    # Convert slides_data to the format expected by the outline_data structure
    # This is needed for batch prompt generation
    simplified_outline = {"title": "Course Presentation", "modules": []}
    # Group slides into temporary modules to match outline format
    temp_module = {"title": "Content Slides", "topics": []}
    for slide in slides_data:
        # Create a topic for each slide
        topic = {
            "title": slide["title"],
            "subtopics": [{"content": slide["content"]}]
        }
        temp_module["topics"].append(topic)
    simplified_outline["modules"].append(temp_module)
    
    # Generate all prompts using the batch approach
    generate_all_image_prompts(simplified_outline, max_slides=len(slides_data))
    elapsed_time = time.time() - start_time
    print(f"Generated all enhanced prompts in {elapsed_time:.1f} seconds")
    
    # Process slides in batches for image generation
    print("\nStep 2: Generating images from enhanced prompts...")
    image_paths_by_title = {}
    
    for i in range(0, len(slides_data), batch_size):
        batch = slides_data[i:i + batch_size]
        print(f"Processing batch {i//batch_size + 1}, slides {i+1} to {min(i+batch_size, len(slides_data))}")
        
        # Prepare prompts and paths for this batch
        prompts_and_paths = []
        for slide in batch:
            # Get enhanced prompt that was pre-generated in batch
            all_slide_text = "\n".join(slide["content"]) if isinstance(slide["content"], list) else slide["content"]
            prompt = get_enhanced_prompt(slide["title"], all_slide_text)
            
            # Create a safe filename
            slide_num = slides_data.index(slide) + 1  # 1-based indexing
            image_filename = f"{slide_num:02d}_slide.png"
            image_path = os.path.join(IMAGE_DIR, image_filename)
            
            # Add to batch
            prompts_and_paths.append((prompt, image_path))
        
        # Generate all images in this batch in parallel
        start_time = time.time()
        all_image_paths = generate_images_parallel(
            prompts_and_paths=prompts_and_paths,
            width=IMAGE_WIDTH,
            height=IMAGE_HEIGHT,
            num_images=1  # One image per slide
        )
        elapsed_time = time.time() - start_time
        
        # Associate images with slides
        successful_images = 0
        for j, slide in enumerate(batch):
            if j < len(all_image_paths):
                image_paths_by_title[slide["title"]] = all_image_paths[j]
                successful_images += 1
        
        print(f"Generated {successful_images}/{len(batch)} images in {elapsed_time:.1f} seconds")
    
    return image_paths_by_title


def add_slides_to_presentation(prs, slides_data, image_paths_by_title):
    """
    Add content slides to the presentation using pre-generated images.
    
    Args:
        prs: PowerPoint presentation object
        slides_data: List of dictionaries with slide information
        image_paths_by_title: Dictionary mapping slide titles to image file paths
        
    Returns:
        None
    """
    global SLIDES_PROCESSED
    
    for slide_data in slides_data:
        title = slide_data["title"]
        content = slide_data.get("content", [])
        color = slide_data.get("color")
        
        # Check if we have a pre-generated image for this slide
        image_path = None
        if image_paths_by_title and title in image_paths_by_title:
            image_path = image_paths_by_title[title]
            if not os.path.exists(image_path):
                print(f"Warning: Image file not found: {image_path}")
                image_path = None
        
        # Create the slide with the image if available
        create_content_slide(prs, title, content, title_color=color, image_path=image_path)
        SLIDES_PROCESSED += 1
        
        # Show progress for large presentations
        if SLIDES_PROCESSED % 10 == 0:
            print(f"Processed {SLIDES_PROCESSED} slides...")

def generate_markdown(outline_data, output_path, course_title="Course Presentation", image_paths_by_title=None):
    """
    Generate a Markdown version of the presentation from the outline data.
    
    Args:
        outline_data: A nested dictionary containing the outline structure
        output_path: Path where the Markdown file will be saved
        course_title: Title of the course for the title slide
        image_paths_by_title: Dictionary mapping slide titles to image paths
        
    Returns:
        True if successful, False otherwise
    """
    try:
        # Get the base directory for relative image paths
        output_dir = os.path.dirname(output_path)
        
        with open(output_path, 'w') as md_file:
            # Write title
            md_file.write(f"# Title: {course_title}\n\n")
            md_file.write(f"*Generated on {time.strftime('%Y-%m-%d')}*\n\n")
            
            # Process modules
            for module_num, module_data in outline_data.items():
                module_title = module_data["title"]
                module_slide_title = f"Module {module_num}: {module_title}"
                md_file.write(f"## {module_title}\n\n")
                
                # Include module slide image if available
                if image_paths_by_title and module_slide_title in image_paths_by_title:
                    image_path = image_paths_by_title[module_slide_title]
                    # Convert to relative path for Markdown
                    rel_path = os.path.relpath(image_path, output_dir)
                    md_file.write(f"![{module_title}]({rel_path})\n\n")
                
                # Process topics in this module
                for topic_num, topic_data in module_data["topics"].items():
                    topic_title = topic_data["title"]
                    topic_slide_title = f"{topic_num}: {topic_title}"
                    md_file.write(f"### {topic_title}\n\n")
                    
                    # Include topic slide image if available
                    if image_paths_by_title and topic_slide_title in image_paths_by_title:
                        image_path = image_paths_by_title[topic_slide_title]
                        # Convert to relative path for Markdown
                        rel_path = os.path.relpath(image_path, output_dir)
                        md_file.write(f"![{topic_title}]({rel_path})\n\n")
                    
                    # Process subtopics in this topic
                    for subtopic_num, subtopic_data in topic_data["subtopics"].items():
                        subtopic_title = subtopic_data["title"]
                        subtopic_slide_title = f"{subtopic_title}"
                        md_file.write(f"#### {subtopic_title}\n\n")
                        
                        # Include subtopic slide image if available
                        if image_paths_by_title and subtopic_slide_title in image_paths_by_title:
                            image_path = image_paths_by_title[subtopic_slide_title]
                            # Convert to relative path for Markdown
                            rel_path = os.path.relpath(image_path, output_dir)
                            md_file.write(f"![{subtopic_title}]({rel_path})\n\n")
                        
                        # Process points in this subtopic
                        for point in subtopic_data["points"]:
                            md_file.write(f"- {point}\n")
                        md_file.write("\n")  # Extra line after points
            
            print(f"Markdown version with images saved to: {output_path}")
            return True
    except Exception as e:
        print(f"Error creating Markdown version: {e}")
        return False
        
def add_slides_to_presentation(prs, slides_data, image_paths_by_title):
    """
    Add content slides to the presentation using pre-generated images.
    
    Args:
        prs: PowerPoint presentation object
        slides_data: List of dictionaries with slide information
        image_paths_by_title: Dictionary mapping slide titles to image file paths
        
    Returns:
        None
    """
    # Now create all slides with the pre-generated images
    print("\nCreating slides with pre-generated images...")
    slides_counter = 1  # Starting with the title slide
    
    for slide_data in slides_data:
        # Create the slide with content
        slide = create_content_slide(
            prs, 
            slide_data["title"], 
            slide_data["content"], 
            slide_data["color"],
            generate_images=False,  # Don't generate images now, we already did
            image_path=image_paths_by_title.get(slide_data["title"])
        )
        slides_counter += 1
        
    print(f"\nGenerated {slides_counter} slides.")
    return slides_counter

def main():
    """Main function to run the presentation generation process."""
    print("=" * 70)
    print("Course PowerPoint Generator".center(70))
    print("=" * 70)
    
    # Determine paths based on current directory
    outline_path = OUTLINE_FILE
    output_pptx_path = OUTPUT_PPTX
    output_markdown_path = OUTPUT_MARKDOWN
    
    if CURRENT_DIR:
        # Use directory-specific paths if available
        outline_path = os.path.join(CURRENT_DIR, OUTLINE_FILE)
        output_pptx_path = os.path.join(CURRENT_DIR, OUTPUT_PPTX)
        output_markdown_path = os.path.join(CURRENT_DIR, OUTPUT_MARKDOWN)
        # Check if we have the outline file in the new directory
        if not os.path.exists(outline_path) and os.path.exists(OUTLINE_FILE):
            # Copy from root to directory
            with open(OUTLINE_FILE, 'r') as src_file:
                content = src_file.read()
            with open(outline_path, 'w') as dst_file:
                dst_file.write(content)
            print(f"Copied outline file to {CURRENT_DIR} directory")
    
    # Check for template file
    template_path = None
    
    # Try each template option
    for template in TEMPLATE_OPTIONS:
        # Check in current directory first
        if CURRENT_DIR and os.path.exists(os.path.join(CURRENT_DIR, template)):
            if template.endswith('.pptx'):
                template_path = os.path.join(CURRENT_DIR, template)
                print(f"\nUsing PowerPoint template from directory: {template_path}")
                break
            elif template.endswith('.potx'):
                print(f"\nFound {template} but python-pptx cannot directly use .potx files.")
                print("Please save your template as .pptx format instead.")
        # Then check in root directory
        elif os.path.exists(template):
            if template.endswith('.pptx'):
                template_path = template
                print(f"\nUsing PowerPoint template: {template}")
                break
            elif template.endswith('.potx'):
                print(f"\nFound {template} but python-pptx cannot directly use .potx files.")
                print("Please save your template as .pptx format instead.")
    
    if template_path is None:
        print("\nNo usable template found. Creating presentation with default styling.")
        print("For custom styling, create a template.pptx file in this directory.")
    
    # Check if the outline file exists
    if not os.path.exists(outline_path):
        print(f"\nError: Course outline file '{outline_path}' not found.")
        print("Please run a04_CREATE_OUTLINE.py first to generate the course outline.")
        return
    
    # Load enhanced notes if available
    notes_file_path = os.path.join(CURRENT_DIR, ENHANCED_NOTES_FILE) if CURRENT_DIR else ENHANCED_NOTES_FILE
    global enhanced_notes
    enhanced_notes = load_enhanced_notes(notes_file_path)
    
    # Debug: Print the keys in enhanced notes to help with matching
    if enhanced_notes:
        print("\nEnhanced notes keys available:")
        for key in enhanced_notes.keys():
            print(f"  - '{key}'")
    else:
        print("\nNo enhanced notes loaded.")
    
    print(f"\nReading course outline from '{outline_path}'...")
    
    try:
        # Parse the outline file
        outline_data, course_title = parse_outline(outline_path)
        
        # Count the number of slides that will be created
        module_count = len(outline_data)
        topic_count = sum(len(module_data["topics"]) for module_data in outline_data.values())
        subtopic_count = sum(
            sum(len(topic_data["subtopics"]) for topic_data in module_data["topics"].values())
            for module_data in outline_data.values()
        )
        
        total_slides = 1 + module_count + topic_count + subtopic_count  # Title slide + content slides
        
        print(f"Found {module_count} modules, {topic_count} topics, and {subtopic_count} subtopics.")
        print(f"Generating PowerPoint presentation with {total_slides} slides...")
        print(f"Using the title: '{course_title}'")
        
        # Prepare slide data for all slides we'll generate
        slides_data = prepare_slides_data(outline_data, max_slides=MAX_SLIDES_TO_PROCESS-1)  # -1 for title slide
        print(f"Preparing {len(slides_data)} slides for generation")
        
        # Generate all images in parallel and collect the paths
        print("Generating slide images...")
        image_paths_by_title = generate_slide_images_parallel(slides_data, batch_size=25)
        
        # Generate the PowerPoint presentation
        print("Creating PowerPoint presentation...")
        prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()
        
        # Generate a dynamic title slide image based on the course title
        title_slide_output_dir = os.path.dirname(output_pptx_path)
        print("\nGenerating title slide image based on course title...")
        title_image_path = generate_title_slide_image(
            course_title, 
            title_slide_output_dir, 
            subtitle_text="Generated Course Presentation"
        )
        
        # Add title slide with the generated image
        print("\nAdding title slide with generated image...")
        create_title_slide(prs, course_title, image_path=title_image_path)
        
        # Add content slides with images
        print("\nCreating slides with pre-generated images...")
        add_slides_to_presentation(prs, slides_data, image_paths_by_title)
        
        # Save the PowerPoint
        prs.save(output_pptx_path)
        
        # Create slide snapshots for Markdown with enhanced layout
        snapshot_dir = os.path.join(os.path.dirname(output_markdown_path), "slide_snapshots")
        os.makedirs(snapshot_dir, exist_ok=True)
        print("\nGenerating enhanced slide snapshots for Markdown export...")
        
        # Find the actual AI-generated images in the output directory
        ai_images_dir = os.path.join(os.path.dirname(output_markdown_path), "slide_images")
        print(f"Looking for AI images in: {ai_images_dir}")
        
        # Create a list of all available image files
        ai_image_files = []
        if os.path.exists(ai_images_dir):
            ai_image_files = sorted([os.path.join(ai_images_dir, f) for f in os.listdir(ai_images_dir) 
                                   if f.endswith('.png') or f.endswith('.jpg')])
            print(f"Found {len(ai_image_files)} AI-generated images")
        
        # Create a mapping of slide indices to image paths
        slide_to_image_map = {}
        for i, slide in enumerate(slides_data):
            if i < len(ai_image_files):
                slide_to_image_map[slide["title"]] = ai_image_files[i]
        
        # Collect all slide data including content
        all_slides_data = []
        for slide in slides_data:
            all_slides_data.append({
                "title": slide["title"],
                "content": slide.get("content", [])
            })
        
        # Generate snapshots using our custom module with corrected image paths
        snapshot_paths = generate_snapshots_for_presentation(
            all_slides_data,
            slide_to_image_map,  # Use our new mapping instead of image_paths_by_title
            snapshot_dir,
            title_image_path=title_image_path,  # Pass the title image path
            course_title=course_title  # Pass the course title
        )
        
        # Generate and save Markdown version with enhanced snapshots
        generate_markdown(outline_data, output_markdown_path, course_title, snapshot_paths)
        
        print("\nGeneration complete!")
        print(f"PowerPoint version: {output_pptx_path}")
        print(f"Markdown version: {output_markdown_path}")
        
        if CURRENT_DIR:
            print(f"Enhanced prompts saved to: {os.path.join(CURRENT_DIR, '06_Enhanced_Prompts.txt')}")
            print(f"Images saved to: {os.path.join(CURRENT_DIR, 'slide_images')}") 
        
        # Add student notes to the PowerPoint
        try:
            # Use correct filename with hyphen instead of underscore
            import importlib.util
            notes_spec = importlib.util.spec_from_file_location("student_notes", "a06-Student_Notes_Student_Handbook.py")
            student_notes = importlib.util.module_from_spec(notes_spec)
            notes_spec.loader.exec_module(student_notes)
            
            success = student_notes.process_presentation_with_notes(outline_data, output_pptx_path)
            if success:
                print(f"\nStudent notes added to presentation '{output_pptx_path}'.")
            else:
                print(f"\nWarning: Student notes could not be added to the presentation.")
        except Exception as e:
            print(f"\nError adding student notes: {str(e)}")
        
        print("="*70)
        
    except Exception as e:
        print(f"\nError occurred: {str(e)}")

if __name__ == "__main__":
    main()
